import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "CargoGuardian"
    subtitle.text = "IoT-Based Transit & Cargo Management System\nContinuous Monitoring & Smart Routing\nTeam: Vaibhav Sharma, Parth Garg, Samarth Sharma, Piyush"

    # Slide 2: Introduction
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Introduction & Problem Addressed"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Transportation of primary resources like coal needs continuous and vigilant monitoring."
    p = tf.add_paragraph()
    p.text = "Current models suffer from manual errors, overloading/underloading issues, and potential frauds."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Need to automate tedious tasks and eliminate human intervention in clearance protocols."
    p.level = 0

    # Slide 3: What is CargoGuardian?
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "What is CargoGuardian?"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "An IoT based project automating cargo transport operations."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Uses load sensors for continuous monitoring of goods mid-transit."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Ensures correct loading and facilitates real-time clearances digitally."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Cross-platform app integration (Flutter) for monitoring from anywhere."
    p.level = 0


    # Slide 4: Core Automated Features
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Core Automated Features"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Zero Human Intervention Clearance: Automatically provides clearance upon confirming correct load via sensors."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Smart Routing: Calculates and detects the shortest possible route to save time and fuel."
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Dynamic Tracking: Live GPS tracks rerouting and monitors the train continuously."
    p.level = 0

    p = tf.add_paragraph()
    p.text = "Fraud Prevention: Instantly detects anomalies (like sudden weight drops) representing potential frauds."
    p.level = 0

    # Slide 5: System Architecture
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "System Architecture & Technologies"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "IoT Layer: Load sensors, GPS Modules, RFID Scanners linked via Blynk."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Backend Layer: Firebase (Realtime Database & Authentication)."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Frontend: Flutter Application (Web, iOS, Android)."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Mapping: Google Maps API utilized for shortest path calculation."
    p.level = 0

    # Slide 6: Future Scope
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Future Scope"

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "Full Automation: Scaling systems to fully manage operations end-to-end."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Truck Integration: Adapting the technology to be used in trucks and the overall road transportation sector."
    p.level = 0
    p = tf.add_paragraph()
    p.text = "AI Analytics: Advanced predictive maintenance and logistical planning based on historical cargo data."
    p.level = 0

    # Slide 7: Thank You
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Thank You"
    subtitle.text = "CargoGuardian - Automating the Future of Logistics"

    prs.save("CargoGuardian_Presentation.pptx")

if __name__ == '__main__':
    create_presentation()
    print("Presentation created successfully as CargoGuardian_Presentation.pptx")
